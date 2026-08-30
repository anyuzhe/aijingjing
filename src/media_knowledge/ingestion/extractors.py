from __future__ import annotations

import hashlib
import html
import io
import json
import math
import re
import shutil
import subprocess
import tempfile
import urllib.error
import urllib.parse
import urllib.request
import wave
from dataclasses import dataclass, field
from html.parser import HTMLParser
from pathlib import Path
from typing import Callable
from uuid import uuid4

from ..models import ContentSegment
from ..product import DesktopSettings, ProductPaths
from ..transcripts import (
    TranscriptQuality as TranscriptV2Quality,
    TranscriptRun as TranscriptV2Run,
    TranscriptSegment as TranscriptV2Segment,
    TranscriptSource as TranscriptV2Source,
    TranscriptSpeaker as TranscriptV2Speaker,
    TranscriptV2,
    TranscriptWord as TranscriptV2Word,
    evaluate_transcript_quality,
    write_transcript,
)
from .audio import AudioPreparationResult, prepare_audio
from .checkpoints import (
    CheckpointIdentity,
    MediaCheckpointStore,
    configuration_hash,
    glossary_version,
)
from .cleanup import TemporaryCleanupRegistry
from .diarization import (
    DiarizationRequest,
    DiarizationResult,
    DiarizationRouter,
    DiarizationSegment,
    DiarizationUnavailable,
    TimedWord,
    build_speaker_cues,
    fuse_words_with_speakers,
)
from .ocr import OCRResult, extract_ocr
from .transcription import (
    TranscriptSegment,
    TranscriptionPlan,
    TranscriptionResult,
    TranscriptionUnavailable,
    transcribe_audio,
    write_transcript_artifacts,
)
from .types import CancelledError, CancellationToken, ExtractionResult
from .vision import MultimodalInterpreter


class MissingExtractorDependency(RuntimeError):
    pass


class PublicDownloadLimitExceeded(ValueError):
    """Raised when a public-platform download crosses the local safety cap."""

    pass


class PublicLiveStreamRejected(ValueError):
    """Raised when a public URL resolves to live or not-yet-final media."""

    pass


class PublicDownloadProtocolRejected(ValueError):
    """Raised when a selected transport cannot be monitored chunk by chunk."""

    pass


@dataclass(slots=True)
class ExtractionContext:
    paths: ProductPaths
    settings: DesktopSettings
    cancellation: CancellationToken
    vision: MultimodalInterpreter | None = None
    progress: Callable[[str], None] | None = None
    cleanup_registry: TemporaryCleanupRegistry | None = field(default=None, repr=False)
    owned_temporary_paths: list[Path] = field(default_factory=list, repr=False)
    cleanup_failures: list[Path] = field(default_factory=list, repr=False)
    temporary_ownership_tokens: dict[Path, str] = field(default_factory=dict, repr=False)

    def message(self, value: str) -> None:
        self.cancellation.check()
        if self.progress:
            self.progress(value)

    def own_temporary_path(self, path: Path) -> Path:
        """Register a cache descendant that this context alone may clean up."""

        if path.is_symlink():
            raise ValueError("临时文件不能是符号链接")
        candidate = path.resolve()
        cache = self.paths.cache.resolve()
        if candidate == cache or cache not in candidate.parents:
            raise ValueError("临时文件必须位于当前产品缓存目录内")
        if candidate not in self.owned_temporary_paths:
            if self.cleanup_registry is not None:
                token = self.cleanup_registry.register(candidate)
                self.temporary_ownership_tokens[candidate] = token
            self.owned_temporary_paths.append(candidate)
        return candidate

    def owns_temporary_path(self, path: Path | None) -> bool:
        if path is None:
            return False
        candidate = path.resolve()
        return any(
            candidate == owned or owned in candidate.parents
            for owned in self.owned_temporary_paths
        )

    def cleanup_temporary_path(self, path: Path, *, attempts: int = 3) -> None:
        candidate = path.resolve()
        if candidate not in self.owned_temporary_paths:
            return
        token = self.temporary_ownership_tokens.get(candidate)
        if self.cleanup_registry is not None and (
            token is None or not self.cleanup_registry.verify(candidate, token)
        ):
            if candidate not in self.cleanup_failures:
                self.cleanup_failures.append(candidate)
            raise OSError("临时清理失败：缓存目录所有权校验未通过，已拒绝删除")
        last_error: OSError | None = None
        for _attempt in range(max(1, attempts)):
            try:
                if candidate.is_dir() and not candidate.is_symlink():
                    shutil.rmtree(candidate)
                else:
                    candidate.unlink(missing_ok=True)
            except OSError as error:
                last_error = error
                continue
            self.owned_temporary_paths.remove(candidate)
            self.temporary_ownership_tokens.pop(candidate, None)
            if candidate in self.cleanup_failures:
                self.cleanup_failures.remove(candidate)
            if self.cleanup_registry is not None and token is not None:
                try:
                    self.cleanup_registry.forget(candidate, token)
                except OSError as error:
                    raise OSError(
                        "临时清理失败：登记更新未完成，将在后续自动重试"
                    ) from error
            return

        if candidate not in self.cleanup_failures:
            self.cleanup_failures.append(candidate)
        raise OSError(
            "临时清理失败：缓存目录已安全登记，已保留待重试记录，将在后续自动重试"
        ) from last_error

    def cleanup_owned_temporary_paths(self) -> None:
        errors: list[OSError] = []
        for path in reversed(self.owned_temporary_paths.copy()):
            try:
                self.cleanup_temporary_path(path)
            except OSError as error:
                errors.append(error)
        if errors:
            raise OSError(
                f"临时清理失败：{len(errors)} 个缓存目录已安全登记，"
                "已保留待重试记录，将在后续自动重试"
            ) from errors[0]


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        while block := handle.read(1024 * 1024):
            digest.update(block)
    return digest.hexdigest()


def safe_stem(value: str) -> str:
    compact = re.sub(r"\s+", " ", value).strip()
    compact = re.sub(r"[/:*?\"<>|\\]", "-", compact).strip(" .-")
    return compact[:96] or "未命名资料"


class TextExtractor:
    suffixes = {".md", ".markdown", ".txt", ".csv", ".tsv", ".json", ".yaml", ".yml", ".py", ".js", ".ts"}

    def extract(self, path: Path, context: ExtractionContext) -> ExtractionResult:
        context.message("正在读取文本")
        text = path.read_text(encoding="utf-8")
        return ExtractionResult(
            title=path.stem,
            media_type="markdown" if path.suffix.casefold() in {".md", ".markdown"} else "text",
            segments=[ContentSegment("text-1", 1, "text", text=text)],
            source_path=path,
            checksum=sha256_file(path),
        )


class PDFExtractor:
    suffixes = {".pdf"}

    def extract(self, path: Path, context: ExtractionContext) -> ExtractionResult:
        try:
            try:
                import pymupdf as fitz  # type: ignore
            except ImportError:
                import fitz  # type: ignore
        except ImportError as exc:
            raise MissingExtractorDependency("PDF 解析组件 PyMuPDF 未安装") from exc
        context.message("正在逐页解析 PDF")
        document = fitz.open(path)
        page_count = len(document)
        segments: list[ContentSegment] = []
        warnings: list[str] = []
        assets: list[Path] = []
        ocr_pages: list[dict[str, object]] = []
        title = str(document.metadata.get("title") or "").strip() or path.stem
        asset_dir = context.paths.assets / "pdf" / safe_stem(path.stem)
        for index, page in enumerate(document):
            context.cancellation.check()
            page_number = index + 1
            text = page.get_text("text").strip()
            description = ""
            should_render = page_number == 1 or len(text) < 100
            image_path: Path | None = None
            if should_render:
                try:
                    asset_dir.mkdir(parents=True, exist_ok=True)
                    image_path = asset_dir / f"page-{page_number:04d}.png"
                    page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False).save(image_path)
                    assets.append(image_path)
                except Exception as exc:
                    warnings.append(f"第 {page_number} 页预览保留失败：{type(exc).__name__}")
            ocr_result: OCRResult | None = None
            if image_path and len(text) < 100:
                context.message(f"正在识别 PDF 第 {page_number} 页扫描内容")
                ocr_result = extract_ocr(
                    image_path,
                    requested_engine=context.settings.ocr_engine,
                    complex_layout=True,
                    allow_paddle=(
                        context.settings.ocr_complex_layout_enabled
                        or context.settings.ocr_engine == "paddleocr"
                    ),
                    low_confidence_threshold=context.settings.ocr_low_confidence_threshold,
                )
                ocr_text = ocr_result.text.strip()
                if ocr_text:
                    if not text:
                        text = ocr_text
                    elif self._normalized_text(ocr_text) not in self._normalized_text(text):
                        text += "\n\n[OCR 补充]\n" + ocr_text
                if ocr_result.fallback_reasons and ocr_result.engine != "paddleocr_ppstructurev3":
                    warnings.append(
                        f"第 {page_number} 页复杂版面 OCR 降级："
                        + "；".join(ocr_result.fallback_reasons)
                    )
            if image_path and context.vision and context.vision.available:
                try:
                    description = context.vision.describe(image_path, context=text[:3000])
                    if ocr_result and (
                        not ocr_result.text
                        or (ocr_result.complex_layout and ocr_result.engine != "paddleocr_ppstructurev3")
                    ):
                        ocr_result.vision_fallback_used = True
                except Exception as exc:
                    warnings.append(f"第 {page_number} 页视觉分析失败：{type(exc).__name__}")
            elif ocr_result and not ocr_result.text:
                ocr_result.fallback_reasons.append("未配置可用的视觉模型，无法执行视觉兜底")
            page_metadata: dict[str, object] = {}
            if ocr_result:
                ocr_metadata = ocr_result.to_dict()
                page_metadata["ocr"] = ocr_metadata
                ocr_pages.append({"page": page_number, **ocr_metadata})
            if not text and not description:
                warnings.append(f"第 {page_number} 页没有提取到文字；可能需要 OCR")
                continue
            segments.append(
                ContentSegment(
                    f"page-{page_number}",
                    page_number,
                    "page",
                    text=text,
                    description=description,
                    location={"page": page_number},
                    asset=str(image_path) if image_path else None,
                    metadata=page_metadata,
                )
            )
        document.close()
        return ExtractionResult(
            title=title,
            media_type="pdf",
            segments=segments,
            source_path=path,
            checksum=sha256_file(path),
            warnings=warnings,
            retained_assets=assets,
            metadata={"page_count": page_count, "ocr": self._ocr_summary(ocr_pages)},
        )

    @staticmethod
    def _normalized_text(value: str) -> str:
        return re.sub(r"\s+", "", value).casefold()

    @staticmethod
    def _ocr_summary(pages: list[dict[str, object]]) -> dict[str, object]:
        scores = [
            float(line["confidence"])
            for page in pages
            for line in page.get("lines", [])  # type: ignore[union-attr]
            if isinstance(line, dict) and isinstance(line.get("confidence"), (int, float))
        ]
        return {
            "pages": pages,
            "ocr_page_count": len(pages),
            "line_count": sum(int(page.get("line_count") or 0) for page in pages),
            "mean_confidence": round(sum(scores) / len(scores), 6) if scores else None,
            "min_confidence": round(min(scores), 6) if scores else None,
            "low_confidence_lines": [
                {"page": page.get("page"), **line}
                for page in pages
                for line in page.get("low_confidence_lines", [])  # type: ignore[union-attr]
                if isinstance(line, dict)
            ],
            "fallback_reasons": list(dict.fromkeys(
                str(reason)
                for page in pages
                for reason in page.get("fallback_reasons", [])  # type: ignore[union-attr]
                if str(reason).strip()
            )),
        }


class DOCXExtractor:
    suffixes = {".docx"}

    def extract(self, path: Path, context: ExtractionContext) -> ExtractionResult:
        try:
            from docx import Document  # type: ignore
        except ImportError as exc:
            raise MissingExtractorDependency("Word 解析组件 python-docx 未安装") from exc
        context.message("正在解析 Word 结构")
        document = Document(path)
        segments: list[ContentSegment] = []
        heading_path: list[str] = []
        sequence = 0
        for paragraph in document.paragraphs:
            context.cancellation.check()
            value = paragraph.text.strip()
            if not value:
                continue
            style = str(paragraph.style.name or "")
            heading = re.match(r"Heading\s+(\d+)", style, re.IGNORECASE)
            if heading:
                level = int(heading.group(1))
                heading_path = heading_path[: level - 1] + [value]
                continue
            sequence += 1
            segments.append(ContentSegment(f"paragraph-{sequence}", sequence, "text", text=value, heading_path=list(heading_path)))
        for table_index, table in enumerate(document.tables, 1):
            rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
            if rows:
                sequence += 1
                segments.append(
                    ContentSegment(
                        f"table-{table_index}", sequence, "table", text="\n".join(rows),
                        description=json.dumps({"rows": len(rows)}, ensure_ascii=False),
                        heading_path=list(heading_path),
                    )
                )
        core_title = str(document.core_properties.title or "").strip()
        return ExtractionResult(
            title=core_title or path.stem,
            media_type="document",
            segments=segments,
            source_path=path,
            checksum=sha256_file(path),
            metadata={"paragraphs": len(document.paragraphs), "tables": len(document.tables)},
        )


class PPTXExtractor:
    suffixes = {".pptx"}

    def extract(self, path: Path, context: ExtractionContext) -> ExtractionResult:
        try:
            from pptx import Presentation  # type: ignore
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
        except ImportError as exc:
            raise MissingExtractorDependency("PPTX 解析组件 python-pptx 未安装") from exc
        context.message("正在逐页解析演示文稿")
        presentation = Presentation(path)
        segments: list[ContentSegment] = []
        assets: list[Path] = []
        warnings: list[str] = []
        asset_dir = context.paths.assets / "slides" / safe_stem(path.stem)
        slide_total = len(presentation.slides)
        for slide_index, slide in enumerate(presentation.slides, 1):
            context.cancellation.check()
            texts: list[str] = []
            descriptions: list[str] = []
            has_visual_structure = False
            for shape_index, shape in enumerate(slide.shapes, 1):
                if getattr(shape, "has_text_frame", False):
                    value = str(shape.text or "").strip()
                    if value:
                        texts.append(value)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_visual_structure = True
                    try:
                        asset_dir.mkdir(parents=True, exist_ok=True)
                        image = shape.image
                        image_path = asset_dir / f"slide-{slide_index:04d}-image-{shape_index:02d}.{image.ext}"
                        image_path.write_bytes(image.blob)
                        assets.append(image_path)
                    except Exception as exc:
                        warnings.append(f"第 {slide_index} 页图片提取失败：{type(exc).__name__}")
                if getattr(shape, "has_table", False):
                    has_visual_structure = True
                if getattr(shape, "has_chart", False):
                    has_visual_structure = True
                    chart_text = self._chart_text(shape)
                    if chart_text:
                        texts.append(chart_text)
            notes = ""
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                notes = str(notes_frame.text or "").strip() if notes_frame else ""
            except (AttributeError, KeyError, ValueError):
                pass
            content = "\n\n".join(texts)
            if notes:
                content += ("\n\n演讲备注：\n" if content else "演讲备注：\n") + notes
            preview_path: Path | None = None
            important_words = ("架构", "流程", "系统", "结果", "曲线", "参数", "对比", "architecture", "flow", "result")
            important = slide_total <= 60 or has_visual_structure or any(word in content.casefold() for word in important_words)
            if important:
                try:
                    asset_dir.mkdir(parents=True, exist_ok=True)
                    preview_path = asset_dir / f"slide-{slide_index:04d}.png"
                    self._render_preview(presentation, slide, preview_path)
                    assets.append(preview_path)
                except Exception as exc:
                    warnings.append(f"第 {slide_index} 页预览保留失败：{type(exc).__name__}")
            if preview_path and context.vision and context.vision.available and (has_visual_structure or len(content) < 180):
                try:
                    descriptions.append(context.vision.describe(preview_path, context=content[:4000]))
                except Exception as exc:
                    warnings.append(f"第 {slide_index} 页视觉分析失败：{type(exc).__name__}")
            if not content and not descriptions:
                warnings.append(f"第 {slide_index} 页没有可索引内容")
                continue
            segments.append(
                ContentSegment(
                    f"slide-{slide_index}", slide_index, "slide", text=content,
                    description="\n\n".join(descriptions), location={"slide": slide_index},
                    heading_path=[texts[0].splitlines()[0][:120]] if texts else [],
                    asset=str(preview_path) if preview_path else None,
                )
            )
        return ExtractionResult(
            title=path.stem,
            media_type="presentation",
            segments=segments,
            source_path=path,
            checksum=sha256_file(path),
            warnings=warnings,
            retained_assets=assets,
            metadata={"slide_count": len(presentation.slides)},
        )

    @staticmethod
    def _chart_text(shape) -> str:
        try:
            chart = shape.chart
            values = []
            if chart.has_title:
                values.append(f"图表标题：{chart.chart_title.text_frame.text}")
            for series in chart.series:
                raw = list(series.values)
                values.append(f"数据系列 {series.name}：{raw[:30]}")
            return "\n".join(values)
        except (AttributeError, TypeError, ValueError):
            return ""

    @staticmethod
    def _render_preview(presentation, slide, destination: Path) -> None:
        try:
            from PIL import Image, ImageDraw, ImageFont  # type: ignore
        except ImportError as exc:
            raise MissingExtractorDependency("PPT 页面预览需要 Pillow") from exc
        width = 1600
        height = max(900, round(width * presentation.slide_height / presentation.slide_width))
        canvas = Image.new("RGB", (width, height), "white")
        draw = ImageDraw.Draw(canvas)
        scale_x = width / presentation.slide_width
        scale_y = height / presentation.slide_height
        font_paths = [
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/STHeiti Medium.ttc",
            "C:/Windows/Fonts/msyh.ttc",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        ]
        font_path = next((value for value in font_paths if Path(value).is_file()), None)

        def font(size: int):
            return ImageFont.truetype(font_path, size) if font_path else ImageFont.load_default()

        for shape in slide.shapes:
            x = round(shape.left * scale_x)
            y = round(shape.top * scale_y)
            w = max(1, round(shape.width * scale_x))
            h = max(1, round(shape.height * scale_y))
            if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                try:
                    picture = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                    picture.thumbnail((w, h))
                    canvas.paste(picture, (x + (w - picture.width) // 2, y + (h - picture.height) // 2))
                    continue
                except (OSError, ValueError):
                    pass
            draw.rounded_rectangle((x, y, min(width - 1, x + w), min(height - 1, y + h)), radius=6, outline="#d6ddd8", width=2)
            text = str(getattr(shape, "text", "") or "").strip()
            if getattr(shape, "has_table", False):
                text = "\n".join(" | ".join(cell.text for cell in row.cells) for row in shape.table.rows)
            if text:
                size = max(16, min(38, round(h / max(3, text.count("\n") + 2))))
                max_chars = max(8, round(w / max(8, size * 0.72)))
                lines = []
                for paragraph in text.splitlines():
                    lines.extend(paragraph[index : index + max_chars] for index in range(0, len(paragraph), max_chars))
                draw.multiline_text((x + 8, y + 6), "\n".join(lines[:18]), fill="#26322d", font=font(size), spacing=5)
        destination.parent.mkdir(parents=True, exist_ok=True)
        canvas.save(destination, optimize=True)


class ImageExtractor:
    suffixes = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp", ".gif"}

    def extract(self, path: Path, context: ExtractionContext) -> ExtractionResult:
        try:
            from PIL import Image  # type: ignore
        except ImportError as exc:
            raise MissingExtractorDependency("图片组件 Pillow 未安装") from exc
        context.message("正在识别图片内容")
        warnings: list[str] = []
        with Image.open(path) as image:
            metadata = {"width": image.width, "height": image.height, "format": image.format}
        ocr_result = extract_ocr(
            path,
            requested_engine=context.settings.ocr_engine,
            allow_paddle=(
                context.settings.ocr_complex_layout_enabled
                or context.settings.ocr_engine == "paddleocr"
            ),
            low_confidence_threshold=context.settings.ocr_low_confidence_threshold,
        )
        ocr = ocr_result.text
        description = ""
        if context.vision and context.vision.available:
            try:
                description = context.vision.describe(path, context=ocr)
                if not ocr or (ocr_result.complex_layout and ocr_result.engine != "paddleocr_ppstructurev3"):
                    ocr_result.vision_fallback_used = True
            except Exception as exc:
                warnings.append(f"视觉分析失败：{type(exc).__name__}")
        elif not ocr:
            ocr_result.fallback_reasons.append("未配置可用的视觉模型，无法执行视觉兜底")
        if ocr_result.fallback_reasons:
            warnings.append("OCR 处理说明：" + "；".join(ocr_result.fallback_reasons))
        if not ocr and not description:
            raise MissingExtractorDependency("图片没有可索引文字，且未配置可用的视觉模型或 OCR")
        ocr_metadata = ocr_result.to_dict()
        metadata["ocr"] = ocr_metadata
        return ExtractionResult(
            title=path.stem,
            media_type="image",
            segments=[ContentSegment(
                "image-1", 1, "image", text=ocr, description=description,
                asset=str(path), metadata={"ocr": ocr_metadata},
            )],
            source_path=path,
            checksum=sha256_file(path),
            warnings=warnings,
            metadata=metadata,
        )


def _ffmpeg_executable() -> str | None:
    executable = shutil.which("ffmpeg")
    if executable:
        return executable
    try:
        import imageio_ffmpeg  # type: ignore

        return imageio_ffmpeg.get_ffmpeg_exe()
    except (ImportError, RuntimeError):
        return None


def _wav_duration(path: Path) -> float:
    try:
        with wave.open(str(path), "rb") as audio:
            rate = audio.getframerate()
            return audio.getnframes() / rate if rate else 0.0
    except (OSError, EOFError, wave.Error):
        return 0.0


def _silence_intervals_ms(prepared: AudioPreparationResult) -> list[tuple[int, int]]:
    duration_ms = max(0, round(prepared.normalized.duration_seconds * 1000))
    cursor = 0
    intervals: list[tuple[int, int]] = []
    for segment in prepared.vad_segments:
        start = max(0, round(segment.start * 1000))
        end = max(start, round(segment.end * 1000))
        if start > cursor:
            intervals.append((cursor, start))
        cursor = max(cursor, end)
    if cursor < duration_ms:
        intervals.append((cursor, duration_ms))
    return intervals


def _media_checkpoint_configuration(settings: DesktopSettings) -> dict[str, object]:
    """Select only settings that can change deterministic media-stage facts."""

    return {
        "pipeline": "audio-video-v2-checkpoints-1",
        "normalization": {
            "sample_rate": 16_000,
            "channels": 1,
            "sample_format": "pcm_s16le",
        },
        "vad": {
            "algorithm": "energy-v1",
            "frame_ms": 30,
            "min_speech_ms": 240,
            "min_silence_ms": 420,
            "padding_ms": 180,
        },
        "asr": {
            "profile": settings.transcription_profile,
            "provider": settings.asr_provider,
            "model": settings.asr_model or settings.whisper_model,
            "model_path": settings.asr_model_path,
            "model_sha256": settings.asr_model_sha256,
            "fallback_model_path": settings.asr_whisper_fallback_model_path,
            "fallback_model_sha256": settings.asr_whisper_fallback_model_sha256,
            "language": settings.transcription_language,
            "engine": settings.transcription_engine,
            "allow_fallback": settings.transcription_allow_cpu_fallback,
            "word_timestamps": settings.word_timestamps,
            "context_terms": list(settings.asr_context_terms),
        },
        "diarization": {
            "enabled": settings.diarization_enabled,
            "provider": settings.diarization_provider,
            "model_path": settings.diarization_model_path,
            "model_sha256": settings.diarization_model_sha256,
            "min_speakers": settings.diarization_min_speakers,
            "max_speakers": settings.diarization_max_speakers,
        },
        "quality_gate": settings.transcript_quality_gate,
    }


def _diarization_from_checkpoint(value: object) -> DiarizationResult | None:
    if not isinstance(value, dict) or value.get("status") != "complete":
        return None
    result = value.get("result")
    if not isinstance(result, dict):
        return None
    raw_segments = result.get("segments")
    segments = [
        DiarizationSegment(
            float(item.get("start", 0.0)),
            float(item.get("end", 0.0)),
            str(item.get("speaker_id") or ""),
            confidence=item.get("confidence"),
            overlap=bool(item.get("overlap", False)),
            metadata=dict(item.get("metadata") or {}),
        )
        for item in raw_segments if isinstance(item, dict)
    ] if isinstance(raw_segments, list) else []
    return DiarizationResult(
        provider_id=str(result.get("provider") or "unknown"),
        model=str(result.get("model") or "unknown"),
        segments=segments,
        fallback_reasons=[str(item) for item in result.get("fallback_reasons", [])],
        warnings=[str(item) for item in result.get("warnings", [])],
        metadata=dict(result.get("metadata") or {}),
    )


def _transcript_v2(
    *,
    source_path: Path,
    source_checksum: str,
    prepared: AudioPreparationResult,
    transcribed: TranscriptionResult,
    settings: DesktopSettings,
    diarization: DiarizationResult | None,
    pipeline_warnings: list[str],
) -> TranscriptV2:
    """Build the immutable/editable Transcript V2 fact representation."""

    run_id = f"asr-run-{uuid4().hex}"
    actual_model_sha256 = settings.asr_model_sha256
    if (
        (transcribed.provider or transcribed.plan.engine) == "mlx-whisper"
        and "qwen3-asr" in (settings.asr_model or "").casefold()
    ):
        actual_model_sha256 = settings.asr_whisper_fallback_model_sha256
    raw_segments = list(transcribed.segments)
    v2_segments: list[TranscriptV2Segment] = []
    if diarization is not None:
        timed_words: list[TimedWord] = []
        for segment in raw_segments:
            if segment.words:
                timed_words.extend(
                    TimedWord(
                        word.start,
                        word.end,
                        word.text,
                        confidence=word.confidence,
                    )
                    for word in segment.words
                    if word.text.strip()
                )
            elif segment.text.strip():
                # Providers without word alignment retain a segment-level timing
                # fact instead of inventing individual word timestamps.
                timed_words.append(TimedWord(
                    segment.start,
                    segment.end,
                    segment.text,
                    confidence=segment.confidence,
                    flags=("segment_level_timing", "speaker_alignment_unavailable"),
                ))
        fused = fuse_words_with_speakers(timed_words, diarization.segments)
        cues = build_speaker_cues(fused)
        for index, cue in enumerate(cues, 1):
            flags = list(cue.flags)
            if cue.overlap:
                flags.append("overlap")
            v2_segments.append(TranscriptV2Segment(
                id=f"{run_id}-seg-{index:04d}",
                ordinal=index - 1,
                start_ms=round(cue.start * 1000),
                end_ms=round(cue.end * 1000),
                speaker_id=cue.speaker_id,
                raw_text=cue.raw_text,
                confidence=cue.confidence,
                flags=tuple(dict.fromkeys(flags)),
                words=tuple(
                    TranscriptV2Word(
                        round(word.start * 1000),
                        round(word.end * 1000),
                        word.text,
                        word.confidence,
                        word.speaker_id,
                        {"overlap": word.overlap},
                    )
                    for word in cue.words
                    if "speaker_alignment_unavailable" not in word.flags
                ),
                metadata={"overlap": cue.overlap},
            ))
    else:
        for index, segment in enumerate(raw_segments, 1):
            v2_segments.append(TranscriptV2Segment(
                id=f"{run_id}-seg-{index:04d}",
                ordinal=index - 1,
                start_ms=round(segment.start * 1000),
                end_ms=round(segment.end * 1000),
                speaker_id=None,
                raw_text=segment.text,
                confidence=segment.confidence,
                flags=(),
                words=tuple(
                    TranscriptV2Word(
                        round(word.start * 1000),
                        round(word.end * 1000),
                        word.text,
                        word.confidence,
                        word.speaker_id,
                    )
                    for word in segment.words
                ),
                metadata={"avg_logprob": segment.avg_logprob},
            ))

    finish_reason = str(transcribed.finish_reason or "stop").strip().casefold() or "stop"
    finish_indicates_truncation = transcribed.truncated or finish_reason in {
        "length", "max_tokens", "token_limit", "truncated", "content_filter",
    }
    if v2_segments:
        final_segment = v2_segments[-1]
        final_segment.metadata.update({
            "finish_reason": finish_reason,
            "transcribed_truncated": bool(transcribed.truncated),
        })
        finish_flags = list(final_segment.flags)
        if finish_indicates_truncation:
            finish_flags.extend(("truncated", "finish_reason_truncated"))
        if finish_reason != "stop":
            finish_flags.append(f"finish_reason_{finish_reason}")
        final_segment.flags = tuple(dict.fromkeys(finish_flags))

    diarization_provider: str | None = None
    if settings.diarization_enabled and settings.diarization_provider != "none":
        diarization_provider = (
            diarization.provider_id
            if diarization is not None
            else f"{settings.diarization_provider}:unavailable"
        )
    transcript = TranscriptV2(
        source=TranscriptV2Source(
            source_path.name,
            source_checksum,
            round(prepared.normalized.duration_seconds * 1000),
            metadata={"audio_probe": prepared.probe.to_dict()},
        ),
        run=TranscriptV2Run(
            id=run_id,
            profile=transcribed.profile,
            provider=transcribed.provider or transcribed.plan.engine,
            model=transcribed.plan.model,
            language=transcribed.language,
            word_timestamps=settings.word_timestamps,
            diarization_provider=diarization_provider,
            context_profile="settings" if settings.asr_context_terms else None,
            fallback={
                "reasons": list(dict.fromkeys([
                    *transcribed.plan.fallback_reasons,
                    *transcribed.fallback_reasons,
                ])),
                "history": list(transcribed.fallback_history),
            } if transcribed.fallback_history or transcribed.fallback_reasons or transcribed.plan.fallback_reasons else None,
            config={
                "device": transcribed.plan.device,
                "compute_type": transcribed.plan.compute_type,
                "requested_provider": settings.asr_provider,
                "requested_model": settings.asr_model or settings.whisper_model,
                "requested_model_sha256": settings.asr_model_sha256,
                "actual_model_sha256": actual_model_sha256,
                "whisper_fallback_model_sha256": (
                    settings.asr_whisper_fallback_model_sha256
                ),
                "language": settings.transcription_language,
                "context_terms": list(settings.asr_context_terms),
                "quality_gate": settings.transcript_quality_gate,
            },
        ),
        speakers=[
            TranscriptV2Speaker(speaker_id)
            for speaker_id in dict.fromkeys(
                item.speaker_id for item in (diarization.segments if diarization else [])
            )
        ],
        segments=v2_segments,
        metadata={
            "vad_segments": [item.to_dict() for item in prepared.vad_segments],
            "audio_probe": prepared.probe.to_dict(),
            "pipeline_warnings": list(dict.fromkeys(pipeline_warnings)),
            "finish_reason": transcribed.finish_reason,
            "truncated": transcribed.truncated,
        },
    )
    report = evaluate_transcript_quality(
        transcript,
        expected_language=(
            settings.transcription_language
            if settings.transcription_language != "auto" else None
        ),
        silence_intervals_ms=_silence_intervals_ms(prepared),
        audio_metrics={
            "decode_ok": prepared.probe.decode_ok,
            "duration_ms": round(prepared.normalized.duration_seconds * 1000),
            "loudness_dbfs": prepared.probe.loudness_dbfs,
            "silence_ratio": prepared.probe.silence_ratio,
            "clipping_ratio": prepared.probe.clipping_ratio,
        },
    )
    extra_warnings = list(dict.fromkeys([
        *prepared.probe.warnings,
        *transcribed.warnings,
        *pipeline_warnings,
        *(diarization.warnings if diarization else []),
    ]))
    status = report.status
    if status == "pass" and extra_warnings:
        status = "review"
    transcript.quality = TranscriptV2Quality(
        status=status,
        warnings=tuple(dict.fromkeys([
            *(issue.message for issue in report.issues),
            *extra_warnings,
        ])),
        metrics={
            **report.metrics,
            "vad_segment_count": len(prepared.vad_segments),
            "decode_ok": prepared.probe.decode_ok,
            "loudness_dbfs": prepared.probe.loudness_dbfs,
            "silence_ratio": prepared.probe.silence_ratio,
            "clipping_ratio": prepared.probe.clipping_ratio,
            "finish_reason": finish_reason,
            "truncated": finish_indicates_truncation,
        },
    )
    return transcript


class AudioVideoExtractor:
    audio_suffixes = {".mp3", ".m4a", ".wav", ".aac", ".flac", ".ogg", ".opus"}
    video_suffixes = {".mp4", ".mov", ".mkv", ".avi", ".webm", ".m4v"}
    suffixes = audio_suffixes | video_suffixes

    def extract(self, path: Path, context: ExtractionContext) -> ExtractionResult:
        ffmpeg = _ffmpeg_executable()
        if not ffmpeg:
            raise MissingExtractorDependency("音视频组件 FFmpeg 未安装或未随应用打包")
        is_video = path.suffix.casefold() in self.video_suffixes
        source_checksum = sha256_file(path)
        checkpoint_config = _media_checkpoint_configuration(context.settings)
        checkpoint_store = MediaCheckpointStore(
            context.paths.cache,
            CheckpointIdentity(
                source_sha256=source_checksum,
                config_hash=configuration_hash(checkpoint_config),
                asr_provider=context.settings.asr_provider,
                asr_model=context.settings.asr_model or context.settings.whisper_model,
                speaker_provider=(
                    context.settings.diarization_provider
                    if context.settings.diarization_enabled else "none"
                ),
                glossary_version=glossary_version(context.settings.asr_context_terms),
                asr_model_sha256=context.settings.asr_model_sha256 or "",
                speaker_model_sha256=(
                    context.settings.diarization_model_sha256 or ""
                ),
            ),
        )
        derived_cache = context.paths.cache / "ingestion-derived"
        derived_cache.mkdir(parents=True, exist_ok=True)
        artifact_root = context.own_temporary_path(
            Path(
                tempfile.mkdtemp(
                    prefix=f"{source_checksum[:12]}-",
                    dir=derived_cache,
                )
            )
        )
        context.message("正在提取音轨")
        try:
            with tempfile.TemporaryDirectory(prefix="ai-jingjing-media-") as temporary:
                prepared = prepare_audio(
                    path,
                    temporary,
                    ffmpeg=ffmpeg,
                    check_cancelled=context.cancellation.check,
                    progress=context.message,
                    checkpoint_store=checkpoint_store,
                )
                audio_path = Path(prepared.normalized.path)
                duration_seconds = prepared.normalized.duration_seconds
                requested_language = {
                    "zh": "Chinese",
                    "en": "English",
                }.get(
                    context.settings.transcription_language,
                    context.settings.transcription_language,
                )
                try:
                    transcribed = transcribe_audio(
                        audio_path,
                        model=(context.settings.asr_model or context.settings.whisper_model),
                        profile=context.settings.transcription_profile,
                        provider=(
                            None
                            if context.settings.asr_provider == "auto"
                            else context.settings.asr_provider
                        ),
                        model_path=context.settings.asr_model_path,
                        whisper_fallback_model_path=(
                            context.settings.asr_whisper_fallback_model_path
                        ),
                        language=None if requested_language == "auto" else requested_language,
                        context_terms=context.settings.asr_context_terms,
                        word_timestamps=context.settings.word_timestamps,
                        preferred_engine=context.settings.transcription_engine,
                        allow_cpu_fallback=context.settings.transcription_allow_cpu_fallback,
                        allow_fallback=context.settings.transcription_allow_cpu_fallback,
                        duration_seconds=duration_seconds,
                        progress=context.message,
                        check_cancelled=context.cancellation.check,
                        checkpoint_store=checkpoint_store,
                    )
                except TranscriptionUnavailable as exc:
                    raise MissingExtractorDependency(str(exc)) from exc
                pipeline_warnings = list(prepared.probe.warnings)
                diarized: DiarizationResult | None = None
                if (
                    context.settings.diarization_enabled
                    and context.settings.diarization_provider != "none"
                ):
                    try:
                        diarized = _diarization_from_checkpoint(
                            checkpoint_store.read_json("diarization.json", "diarization")
                        )
                    except (TypeError, ValueError, OverflowError):
                        diarized = None
                    if diarized is not None:
                        context.message("已复用说话人分段检查点")
                    else:
                        context.message("正在区分并对齐说话人")
                        try:
                            diarized = DiarizationRouter().diarize(
                                DiarizationRequest(
                                    audio_path=audio_path,
                                    model_path=(
                                        Path(context.settings.diarization_model_path)
                                        if context.settings.diarization_model_path else None
                                    ),
                                    preferred_provider=context.settings.diarization_provider,
                                    min_speakers=context.settings.diarization_min_speakers,
                                    max_speakers=context.settings.diarization_max_speakers,
                                    allow_fallback=True,
                                ),
                                progress=context.message,
                                check_cancelled=context.cancellation.check,
                            )
                            checkpoint_store.write_json(
                                "diarization.json",
                                "diarization",
                                {"status": "complete", "result": diarized.to_dict()},
                                runtime={"speaker_provider": diarized.provider_id},
                            )
                        except CancelledError:
                            raise
                        except (DiarizationUnavailable, OSError, RuntimeError, ValueError) as exc:
                            pipeline_warnings.append(f"说话人识别需要复核：{exc}")
                            checkpoint_store.write_json(
                                "diarization.json",
                                "diarization",
                                {"status": "unavailable", "reason": str(exc)},
                            )
                else:
                    checkpoint_store.write_json(
                        "diarization.json",
                        "diarization",
                        {"status": "disabled", "result": None},
                    )

                transcript_v2 = _transcript_v2(
                    source_path=path,
                    source_checksum=source_checksum,
                    prepared=prepared,
                    transcribed=transcribed,
                    settings=context.settings,
                    diarization=diarized,
                    pipeline_warnings=pipeline_warnings,
                )
                checkpoint_runtime = {
                    "asr_provider": transcribed.provider or transcribed.plan.engine,
                    "asr_model": transcribed.plan.model,
                    "asr_model_sha256": (
                        context.settings.asr_whisper_fallback_model_sha256
                        if (
                            (transcribed.provider or transcribed.plan.engine)
                            == "mlx-whisper"
                            and "qwen3-asr"
                            in (context.settings.asr_model or "").casefold()
                        )
                        else context.settings.asr_model_sha256
                    ),
                    "speaker_provider": (
                        diarized.provider_id if diarized else checkpoint_store.identity.speaker_provider
                    ),
                }
                transcript_v2.metadata["checkpoint"] = checkpoint_store.checkpoint_metadata(
                    "transcript_v2", runtime=checkpoint_runtime
                )
                write_transcript(
                    transcript_v2,
                    checkpoint_store.path("transcript-v2.json"),
                )
                checkpoint_store.record_file(
                    "transcript-v2.json", "transcript_v2", runtime=checkpoint_runtime
                )
                checkpoint_store.write_json(
                    "quality.json",
                    "quality",
                    transcript_v2.quality.to_dict(),
                    runtime=checkpoint_runtime,
                )
                segments: list[ContentSegment] = []
                speaker_names = {
                    speaker.id: speaker.display_name or speaker.id
                    for speaker in transcript_v2.speakers
                }
                for index, item in enumerate(transcript_v2.segments, 1):
                    context.cancellation.check()
                    if not item.effective_text:
                        continue
                    start = item.start_ms / 1000.0
                    end = item.end_ms / 1000.0
                    segments.append(ContentSegment(
                        item.id, start, "speech", text=item.effective_text,
                        location={
                            "timestamp_start": start,
                            "timestamp_end": end,
                            "speaker_id": item.speaker_id,
                        },
                        metadata={
                            "language": transcribed.language,
                            "engine": transcribed.plan.engine,
                            "provider": transcribed.provider,
                            "model": transcribed.plan.model,
                            "confidence": item.confidence,
                            "speaker_id": item.speaker_id,
                            "speaker_name": speaker_names.get(item.speaker_id or ""),
                            "overlap": "overlap" in item.flags,
                            "asr_run_id": transcript_v2.run.id,
                            "quality_status": transcript_v2.quality.status,
                            "quality_flags": list(item.flags),
                            "raw_text": item.raw_text,
                        },
                    ))
                transcript_basename = f"{safe_stem(path.stem)}-{source_checksum[:10]}"
                transcript_artifacts = write_transcript_artifacts(
                    transcribed,
                    artifact_root / "transcript",
                    transcript_basename,
                    source_name=path.name,
                )
                transcript_v2_path = write_transcript(
                    transcript_v2,
                    artifact_root / "transcript" / f"{transcript_basename}.v2.json",
                )
                transcribed.artifacts["v2"] = str(transcript_v2_path)
                transcript_path = transcript_artifacts["txt"]
                assets: list[Path] = []
                warnings = list(dict.fromkeys([
                    *prepared.probe.warnings,
                    *transcribed.plan.fallback_reasons,
                    *transcribed.fallback_reasons,
                    *transcribed.warnings,
                    *pipeline_warnings,
                    *transcript_v2.quality.warnings,
                ]))
                if is_video and context.vision and context.vision.available:
                    context.message("正在抽取视频关键帧")
                    frame_dir = artifact_root / "frames"
                    frame_dir.mkdir(parents=True, exist_ok=True)
                    subprocess.run(
                        [
                            ffmpeg, "-y", "-i", str(path), "-vf", "fps=1/60,scale=1280:-2",
                            "-frames:v", "8", str(frame_dir / "frame-%03d.jpg"),
                        ],
                        capture_output=True,
                        timeout=30 * 60,
                    )
                    for frame_index, frame in enumerate(sorted(frame_dir.glob("frame-*.jpg")), 1):
                        try:
                            assets.append(frame)
                            description = context.vision.describe(frame, context="视频关键帧")
                            segments.append(
                                ContentSegment(
                                    f"frame-{frame_index}", frame_index * 60, "image",
                                    description=description,
                                    location={"timestamp_start": float((frame_index - 1) * 60)},
                                    asset=str(frame),
                                )
                            )
                        except Exception as exc:
                            warnings.append(f"关键帧 {frame_index} 分析失败：{type(exc).__name__}")
            if not segments:
                raise RuntimeError("音视频没有提取到可索引内容")
        except BaseException as error:
            try:
                context.cleanup_temporary_path(artifact_root)
            except OSError as cleanup_error:
                error.add_note(f"附加诊断：{cleanup_error}")
            raise
        return ExtractionResult(
            title=path.stem,
            media_type="video" if is_video else "audio",
            segments=segments,
            source_path=path,
            checksum=source_checksum,
            warnings=warnings,
            retained_assets=assets,
            transcript_path=transcript_path,
            transcript_data=transcript_v2.to_dict(),
            metadata={
                "transcription": {
                    **transcribed.metadata(),
                    "v2_run_id": transcript_v2.run.id,
                    "quality": transcript_v2.quality.to_dict(),
                    "audio_preparation": prepared.to_dict(),
                    "diarization": diarized.to_dict() if diarized else None,
                },
            },
        )


_BROWSER_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
        "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/140.0 Safari/537.36"
    ),
    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8",
    "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.7",
}


class _ReadableHTML(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.title = ""
        self._in_title = False
        self._ignored = 0
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs) -> None:
        attributes = dict(attrs)
        if tag == "meta" and attributes.get("property") in {"og:title", "twitter:title"}:
            self.title = str(attributes.get("content") or self.title)
        if tag in {"script", "style", "noscript", "svg"}:
            self._ignored += 1
        if tag == "title":
            self._in_title = True
        if tag in {"p", "div", "section", "article", "h1", "h2", "h3", "li", "br"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript", "svg"} and self._ignored:
            self._ignored -= 1
        if tag == "title":
            self._in_title = False

    def handle_data(self, data: str) -> None:
        if self._ignored:
            return
        if self._in_title:
            self.title += data
        else:
            self.parts.append(data)

    def text(self) -> str:
        value = html.unescape(" ".join(self.parts))
        value = re.sub(r"[ \t\r\f\v]+", " ", value)
        value = re.sub(r"\n\s*\n+", "\n\n", value)
        return value.strip()


class _WeixinArticleHTML(HTMLParser):
    """Extract the real article body from a Weixin public-account page."""

    _void_tags = {
        "area", "base", "br", "col", "embed", "hr", "img", "input",
        "link", "meta", "param", "source", "track", "wbr",
    }
    _block_tags = {"p", "div", "section", "article", "h1", "h2", "h3", "h4", "li", "blockquote", "br", "hr"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.title = ""
        self.parts: list[str] = []
        self._content_depth = 0
        self._ignored = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        attributes = dict(attrs)
        if tag == "meta" and attributes.get("property") in {"og:title", "twitter:title"}:
            self.title = str(attributes.get("content") or self.title)
        entering_content = attributes.get("id") == "js_content"
        if entering_content:
            self._content_depth = 1
        elif self._content_depth and tag not in self._void_tags:
            self._content_depth += 1
        if self._content_depth and tag in {"script", "style", "noscript", "svg"}:
            self._ignored += 1
        if self._content_depth and tag in self._block_tags:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if not self._content_depth:
            return
        if tag in {"script", "style", "noscript", "svg"} and self._ignored:
            self._ignored -= 1
        if tag in self._block_tags:
            self.parts.append("\n")
        if tag not in self._void_tags:
            self._content_depth -= 1

    def handle_data(self, data: str) -> None:
        if self._content_depth and not self._ignored:
            self.parts.append(data)

    def text(self) -> str:
        value = html.unescape(" ".join(self.parts))
        value = re.sub(r"[ \t\r\f\v]+", " ", value)
        value = re.sub(r"\n\s*\n+", "\n\n", value)
        return value.strip()


def _download_web_document(url: str, context: ExtractionContext) -> tuple[bytes, str, str]:
    context.message("正在下载网页快照")
    request = urllib.request.Request(url, headers=_BROWSER_HEADERS)
    with urllib.request.urlopen(request, timeout=45) as response:
        raw = response.read(15 * 1024 * 1024 + 1)
        if len(raw) > 15 * 1024 * 1024:
            raise ValueError("网页超过 15MB 安全限制")
        charset = response.headers.get_content_charset() or "utf-8"
        document = raw.decode(charset, errors="replace")
        final_url = response.geturl() if hasattr(response, "geturl") else url
    return raw, document, final_url


def _raise_for_weixin_challenge(final_url: str, document: str) -> None:
    path = urllib.parse.urlsplit(final_url).path.casefold()
    challenge_markers = (
        "当前环境异常，完成验证后即可继续访问",
        "wappoc_appmsgcaptcha",
        "访问过于频繁，请用微信扫描二维码进行访问",
    )
    if path.startswith("/mp/wappoc_appmsgcaptcha") or any(marker in document for marker in challenge_markers):
        raise RuntimeError("微信返回了访问验证页，没有取得文章正文；请稍后重试或在浏览器中打开后导出为 PDF 再导入")


class WebExtractor:
    def extract(self, url: str, context: ExtractionContext) -> ExtractionResult:
        raw, document, _final_url = _download_web_document(url, context)
        parser = _ReadableHTML()
        parser.feed(document)
        text = parser.text()
        if not text:
            raise RuntimeError("网页没有提取到正文；可能需要登录或动态浏览器渲染")
        return ExtractionResult(
            title=parser.title.strip() or url,
            media_type="web",
            segments=[ContentSegment("web-1", 1, "text", text=text)],
            original_uri=url,
            checksum=hashlib.sha256(raw).hexdigest(),
            metadata={"snapshot_html": document},
        )


class WeixinArticleExtractor:
    """Read public Weixin articles without treating anti-bot pages as knowledge."""

    @classmethod
    def supports(cls, url: str) -> bool:
        parsed = urllib.parse.urlsplit(url)
        host = (parsed.hostname or "").casefold()
        return parsed.scheme.casefold() in {"http", "https"} and host == "mp.weixin.qq.com" and parsed.path.startswith("/s/")

    def extract(self, url: str, context: ExtractionContext) -> ExtractionResult:
        raw, document, final_url = _download_web_document(url, context)
        _raise_for_weixin_challenge(final_url, document)
        parser = _WeixinArticleHTML()
        parser.feed(document)
        text = parser.text()
        if len(text) < 80:
            raise RuntimeError("微信公众号页面没有提取到足够的文章正文；页面可能已删除、需要验证或限制外部访问")
        return ExtractionResult(
            title=parser.title.strip() or url,
            media_type="web",
            segments=[ContentSegment("web-1", 1, "text", text=text)],
            original_uri=url,
            checksum=hashlib.sha256(raw).hexdigest(),
            metadata={
                "snapshot_html": document,
                "platform": "weixin_public_account",
                "content_scope": "full_source",
            },
        )


class DirectMediaURLExtractor:
    """Download a public direct audio/video URL, then run the local media pipeline."""

    media_suffixes = AudioVideoExtractor.suffixes
    max_download_bytes = 2 * 1024 * 1024 * 1024
    content_type_suffixes = {
        "audio/aac": ".aac",
        "audio/flac": ".flac",
        "audio/m4a": ".m4a",
        "audio/mp4": ".m4a",
        "audio/mpeg": ".mp3",
        "audio/ogg": ".ogg",
        "audio/opus": ".opus",
        "audio/wav": ".wav",
        "audio/x-wav": ".wav",
        "video/mp4": ".mp4",
        "video/quicktime": ".mov",
        "video/webm": ".webm",
        "video/x-m4v": ".m4v",
        "video/x-matroska": ".mkv",
        "video/x-msvideo": ".avi",
    }

    @classmethod
    def supports(cls, url: str) -> bool:
        parsed = urllib.parse.urlsplit(url)
        return parsed.scheme.casefold() in {"http", "https"} and Path(parsed.path).suffix.casefold() in cls.media_suffixes

    def extract(self, url: str, context: ExtractionContext) -> ExtractionResult:
        return self.extract_remote(url, context, original_uri=url)

    def extract_remote(
        self,
        media_url: str,
        context: ExtractionContext,
        *,
        original_uri: str,
        title: str | None = None,
        metadata: dict[str, object] | None = None,
        request_headers: dict[str, str] | None = None,
    ) -> ExtractionResult:
        parsed = urllib.parse.urlsplit(media_url)
        if parsed.scheme.casefold() not in {"http", "https"}:
            raise ValueError("远程音视频地址必须使用 HTTP 或 HTTPS")
        headers = {"User-Agent": "Mozilla/5.0 AI-Jingjing/1.0", **(request_headers or {})}
        request = urllib.request.Request(media_url, headers=headers)
        context.message("正在下载远程音视频")
        identity = hashlib.sha256(original_uri.encode("utf-8")).hexdigest()[:16]
        cache = context.paths.cache / "remote-media"
        cache.mkdir(parents=True, exist_ok=True)
        directory = context.own_temporary_path(
            Path(tempfile.mkdtemp(prefix=f"{identity}-", dir=cache))
        )
        try:
            with urllib.request.urlopen(request, timeout=90) as response:
                content_type = str(response.headers.get_content_type() or "").casefold()
                declared_length = response.headers.get("Content-Length")
                if declared_length:
                    try:
                        if int(declared_length) > self.max_download_bytes:
                            raise ValueError("远程音视频超过 2GB 安全限制，请先下载后再导入")
                    except ValueError as exc:
                        if "超过" in str(exc):
                            raise
                suffix = Path(parsed.path).suffix.casefold()
                if suffix not in self.media_suffixes:
                    suffix = self.content_type_suffixes.get(content_type, "")
                if not (content_type.startswith(("audio/", "video/")) or suffix in self.media_suffixes):
                    raise ValueError("该地址不是可直接下载的音视频链接")
                if suffix not in self.media_suffixes:
                    suffix = ".mp4" if content_type.startswith("video/") else ".m4a"
                source_stem = safe_stem(title or Path(parsed.path).stem or "远程音视频")
                destination = directory / f"{source_stem}{suffix}"
                temporary = destination.with_suffix(destination.suffix + ".part")
                downloaded = 0
                with temporary.open("wb") as handle:
                    while block := response.read(1024 * 1024):
                        context.cancellation.check()
                        downloaded += len(block)
                        if downloaded > self.max_download_bytes:
                            raise ValueError("远程音视频超过 2GB 安全限制，请先下载后再导入")
                        handle.write(block)
                temporary.replace(destination)
            extracted = AudioVideoExtractor().extract(destination, context)
            extracted.source_path = destination
            extracted.title = title or extracted.title
            extracted.original_uri = original_uri
            extracted.metadata.update(
                {
                    "remote_media": True,
                    "remote_content_type": content_type,
                    "remote_download_bytes": downloaded,
                    "temporary_source_owned_by": "ExtractionContext",
                    "temporary_source_identity": identity,
                    "source_media_bytes": destination.stat().st_size,
                    **(metadata or {}),
                }
            )
            return extracted
        except BaseException as error:
            try:
                context.cleanup_temporary_path(directory)
            except OSError as cleanup_error:
                error.add_note(f"附加诊断：{cleanup_error}")
            raise


class _QuietYTDLPLogger:
    """Discard third-party logs so signed CDN URLs or credentials cannot leak."""

    def debug(self, _message: str) -> None:
        pass

    def warning(self, _message: str) -> None:
        pass

    def error(self, _message: str) -> None:
        pass


def _youtube_dl_class():
    try:
        from yt_dlp import YoutubeDL  # type: ignore
    except ImportError as exc:
        raise MissingExtractorDependency(
            "公开视频平台连接器 yt-dlp 未安装；请安装 media 可选组件，"
            "或先把公开媒体保存为本地文件后再导入"
        ) from exc
    return YoutubeDL


def _subtitle_seconds(value: str) -> float | None:
    match = re.fullmatch(r"(?:(\d+):)?(\d{1,2}):(\d{2})[,.](\d{3})", value.strip())
    if not match:
        return None
    hours = int(match.group(1) or 0)
    minutes = int(match.group(2))
    seconds = int(match.group(3))
    milliseconds = int(match.group(4))
    if minutes >= 60 or seconds >= 60:
        return None
    return hours * 3600 + minutes * 60 + seconds + milliseconds / 1000


def _parse_subtitle_file(path: Path) -> list[TranscriptSegment]:
    """Parse VTT/SRT without executing styles, HTML, or embedded metadata."""

    raw = path.read_text(encoding="utf-8-sig", errors="replace")
    blocks = re.split(r"\n\s*\n", raw.replace("\r\n", "\n").replace("\r", "\n"))
    segments: list[TranscriptSegment] = []
    for block in blocks:
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        if not lines or lines[0].upper().startswith(("WEBVTT", "NOTE", "STYLE", "REGION")):
            continue
        timing_index = next((index for index, line in enumerate(lines[:3]) if "-->" in line), None)
        if timing_index is None:
            continue
        timing = lines[timing_index].split("-->", 1)
        start = _subtitle_seconds(timing[0].strip())
        end_token = timing[1].strip().split()[0] if timing[1].strip() else ""
        end = _subtitle_seconds(end_token)
        if start is None or end is None:
            continue
        text = "\n".join(lines[timing_index + 1 :])
        text = html.unescape(re.sub(r"<[^>]{0,500}>", "", text)).replace("\x00", "").strip()
        if segments and segments[-1].text == text and start <= segments[-1].end + 0.25:
            previous = segments[-1]
            segments[-1] = TranscriptSegment(previous.start, max(previous.end, end), text)
        else:
            segments.append(TranscriptSegment(start, end, text))
    return segments


class PublicPlatformVideoExtractor:
    """Public-only yt-dlp connector with subtitles-first ingestion.

    The connector never asks yt-dlp for browser cookies, never supplies a cookie
    file, and explicitly disables proxies. Private/authenticated media is not
    retried with user credentials.
    """

    host_platforms = {
        "youtube.com": "youtube",
        "www.youtube.com": "youtube",
        "m.youtube.com": "youtube",
        "youtu.be": "youtube",
        "bilibili.com": "bilibili",
        "www.bilibili.com": "bilibili",
        "m.bilibili.com": "bilibili",
        "b23.tv": "bilibili",
        "douyin.com": "douyin",
        "www.douyin.com": "douyin",
        "v.douyin.com": "douyin",
        "xiaohongshu.com": "xiaohongshu",
        "www.xiaohongshu.com": "xiaohongshu",
        "xhslink.com": "xiaohongshu",
        "www.xhslink.com": "xiaohongshu",
        "x.com": "x",
        "www.x.com": "x",
        "twitter.com": "x",
        "www.twitter.com": "x",
        "mobile.twitter.com": "x",
    }
    max_download_bytes = 2 * 1024 * 1024 * 1024

    @classmethod
    def supports(cls, url: str) -> bool:
        try:
            parsed = urllib.parse.urlsplit(url)
            return (
                parsed.scheme.casefold() in {"http", "https"}
                and parsed.hostname is not None
                and parsed.hostname.casefold() in cls.host_platforms
                and parsed.username is None
                and parsed.password is None
                and parsed.port in {None, 80, 443}
            )
        except ValueError:
            return False

    @staticmethod
    def _safe_info(raw: object) -> dict[str, object]:
        info = raw if isinstance(raw, dict) else {}
        entries = info.get("entries")
        if isinstance(entries, list):
            entries = [item for item in entries if isinstance(item, dict)]
            if len(entries) != 1:
                raise ValueError("一次只能导入一个公开视频，暂不支持播放列表批量下载")
            info = entries[0]
        return info

    @staticmethod
    def _subtitle_choice(info: dict[str, object]) -> tuple[str, bool] | None:
        priorities = ("zh-Hans", "zh-CN", "zh-TW", "zh", "en")
        for automatic, key in ((False, "subtitles"), (True, "automatic_captions")):
            values = info.get(key)
            if not isinstance(values, dict) or not values:
                continue
            languages = [str(item) for item in values]
            language = next((item for item in priorities if item in languages), None)
            if language is None:
                language = sorted(languages)[0]
            return language, automatic
        return None

    @staticmethod
    def _byte_count(value: object) -> int:
        if isinstance(value, bool) or value is None:
            return 0
        try:
            number = float(value)
        except (TypeError, ValueError, OverflowError):
            return 0
        if not math.isfinite(number) or number <= 0:
            return 0
        return int(number)

    @classmethod
    def _declared_download_size(cls, info: dict[str, object]) -> int:
        """Return the strongest size declaration exposed by yt-dlp metadata."""

        top_level = max(
            cls._byte_count(info.get("filesize")),
            cls._byte_count(info.get("filesize_approx")),
        )
        requested = info.get("requested_formats")
        requested_total = 0
        if isinstance(requested, list):
            for item in requested:
                if not isinstance(item, dict):
                    continue
                requested_total += max(
                    cls._byte_count(item.get("filesize")),
                    cls._byte_count(item.get("filesize_approx")),
                )
        return max(top_level, requested_total)

    @staticmethod
    def _downloaded_directory_bytes(directory: Path) -> int:
        """Measure regular downloaded files without following symlinks."""

        total = 0
        try:
            candidates = directory.rglob("*")
            for candidate in candidates:
                try:
                    if candidate.is_symlink() or not candidate.is_file():
                        continue
                    total += candidate.stat().st_size
                except OSError:
                    # yt-dlp may rename fragments while a hook is running.
                    continue
        except OSError:
            return total
        return total

    @staticmethod
    def _limit_error() -> PublicDownloadLimitExceeded:
        return PublicDownloadLimitExceeded(
            "公开视频下载超过 2GB 安全上限，已停止并清理临时文件；"
            "请先下载并裁剪后再导入"
        )

    def _enforce_declared_limit(self, info: dict[str, object]) -> None:
        if self._declared_download_size(info) > self.max_download_bytes:
            raise self._limit_error()

    @staticmethod
    def _selected_format_records(info: dict[str, object]) -> list[dict[str, object]]:
        records = [info]
        for key in ("requested_formats", "requested_downloads", "selected_formats"):
            value = info.get(key)
            if isinstance(value, dict):
                records.append(value)
            elif isinstance(value, list):
                records.extend(item for item in value if isinstance(item, dict))
        return records

    @classmethod
    def _is_live_selection(cls, info: dict[str, object]) -> bool:
        rejected_statuses = {"is_live", "live", "is_upcoming", "upcoming", "post_live"}
        for record in cls._selected_format_records(info):
            if record.get("is_live") is True:
                return True
            if str(record.get("live_status") or "").strip().casefold() in rejected_statuses:
                return True
        return False

    @staticmethod
    def _live_error() -> PublicLiveStreamRejected:
        return PublicLiveStreamRejected(
            "为保证下载上限和内容完整性，AI静静不导入直播、即将直播或仍在生成回放的内容；"
            "请等待回放完成，或保存为本地文件后导入"
        )

    @classmethod
    def _enforce_not_live(cls, info: dict[str, object]) -> None:
        if cls._is_live_selection(info):
            raise cls._live_error()

    @staticmethod
    def _protocol_error(protocol: str) -> PublicDownloadProtocolRejected:
        safe_protocol = re.sub(r"[^a-z0-9_+.-]", "", protocol.casefold())[:48] or "unknown"
        return PublicDownloadProtocolRejected(
            f"公开视频使用了不可安全监控的传输协议（{safe_protocol}）；"
            "当前仅支持可逐块限制的 HTTP、HTTPS 和原生 DASH 分片。"
            "HLS（m3u8）本版本不会在线下载，可先合法保存为本地文件后导入"
        )

    @classmethod
    def _enforce_safe_protocols(cls, info: dict[str, object]) -> None:
        allowed = {
            "http",
            "https",
            "http_dash_segments",
            "http_dash_segments_generator",
        }
        for record in cls._selected_format_records(info):
            raw = str(record.get("protocol") or "").strip().casefold()
            if not raw:
                continue
            protocols = [part for part in raw.split("+") if part]
            rejected = next((part for part in protocols if part not in allowed), None)
            if rejected:
                raise cls._protocol_error(rejected)

    def _match_filter(self, info: dict[str, object], *, incomplete: bool = False) -> None:
        del incomplete
        safe_info = info if isinstance(info, dict) else {}
        self._enforce_not_live(safe_info)
        self._enforce_safe_protocols(safe_info)
        return None

    def _progress_hook(self, context: ExtractionContext, directory: Path):
        state = {"bucket": -1}

        def hook(data: dict[str, object]) -> None:
            context.cancellation.check()
            reported_downloaded = self._byte_count(data.get("downloaded_bytes"))
            reported_total = max(
                self._byte_count(data.get("total_bytes")),
                self._byte_count(data.get("total_bytes_estimate")),
            )
            on_disk = self._downloaded_directory_bytes(directory)
            if max(reported_downloaded, reported_total, on_disk) > self.max_download_bytes:
                raise self._limit_error()
            status = str(data.get("status") or "")
            if status == "downloading":
                percent = re.sub(r"\x1b\[[0-9;]*m", "", str(data.get("_percent_str") or "")).strip()
                try:
                    bucket = int(float(percent.rstrip("%")) // 5)
                except (TypeError, ValueError):
                    bucket = 0
                if bucket != state["bucket"]:
                    state["bucket"] = bucket
                    context.message(f"正在下载公开视频{f'（{percent}）' if percent else ''}")
            elif status == "finished":
                context.message("公开视频下载完成，正在校验")
        return hook

    def _options(self, directory: Path, context: ExtractionContext) -> dict[str, object]:
        return {
            "quiet": True,
            "no_warnings": True,
            "logger": _QuietYTDLPLogger(),
            "noplaylist": True,
            "restrictfilenames": True,
            "overwrites": True,
            "continuedl": True,
            "proxy": "",
            "cookiefile": None,
            "cookiesfrombrowser": None,
            "usenetrc": False,
            "netrc_location": None,
            "socket_timeout": 30,
            "retries": 2,
            "fragment_retries": 2,
            "max_filesize": self.max_download_bytes,
            "external_downloader": {
                "default": "native",
                "dash": "native",
            },
            "match_filter": self._match_filter,
            "outtmpl": str(directory / "%(id)s.%(ext)s"),
            "progress_hooks": [self._progress_hook(context, directory)],
        }

    @staticmethod
    def _control_exception(
        error: BaseException,
    ) -> (
        CancelledError
        | PublicDownloadLimitExceeded
        | PublicLiveStreamRejected
        | PublicDownloadProtocolRejected
        | None
    ):
        """Recover control-flow exceptions that yt-dlp may wrap internally."""

        pending: list[BaseException] = [error]
        visited: set[int] = set()
        while pending:
            current = pending.pop()
            if id(current) in visited:
                continue
            visited.add(id(current))
            if isinstance(
                current,
                (
                    CancelledError,
                    PublicDownloadLimitExceeded,
                    PublicLiveStreamRejected,
                    PublicDownloadProtocolRejected,
                ),
            ):
                return current
            for related in (current.__cause__, current.__context__):
                if isinstance(related, BaseException):
                    pending.append(related)
            exc_info = getattr(current, "exc_info", None)
            if (
                isinstance(exc_info, tuple)
                and len(exc_info) >= 2
                and isinstance(exc_info[1], BaseException)
            ):
                pending.append(exc_info[1])
        return None

    def _run_ytdlp(self, url: str, options: dict[str, object], *, download: bool) -> dict[str, object]:
        YoutubeDL = _youtube_dl_class()
        try:
            with YoutubeDL(options) as ydl:
                return self._safe_info(ydl.extract_info(url, download=download))
        except (
            CancelledError,
            PublicDownloadLimitExceeded,
            PublicLiveStreamRejected,
            PublicDownloadProtocolRejected,
        ):
            raise
        except Exception as exc:
            control = self._control_exception(exc)
            if control is not None:
                raise control
            # Do not echo yt-dlp's exception because it can contain signed media
            # URLs. The recovery path is actionable and credential-free.
            raise RuntimeError(
                "无法取得该平台的公开字幕或媒体。AI静静没有读取浏览器 Cookie、"
                "没有使用代理，也不会尝试绕过登录/地区/权限限制；可先公开下载为本地文件后导入。"
            ) from exc

    @staticmethod
    def _duration(info: dict[str, object], segments: list[TranscriptSegment]) -> float:
        try:
            duration = float(info.get("duration") or 0.0)
        except (TypeError, ValueError, OverflowError):
            duration = 0.0
        if duration <= 0 and segments:
            duration = max(item.end for item in segments)
        return max(0.0, duration)

    def _subtitle_result(
        self,
        *,
        url: str,
        platform: str,
        info: dict[str, object],
        sidecar: Path,
        language: str,
        automatic: bool,
        context: ExtractionContext,
    ) -> ExtractionResult:
        segments = _parse_subtitle_file(sidecar)
        if not any(item.text for item in segments):
            raise RuntimeError("平台字幕文件为空或格式无法识别，将改用音视频转写")
        title = str(info.get("title") or info.get("id") or platform).strip()
        duration = self._duration(info, segments)
        transcript = TranscriptionResult(
            plan=TranscriptionPlan("platform-subtitle", "public-sidecar", "source", "none"),
            language=language,
            duration_seconds=duration,
            segments=segments,
        )
        from .quality import evaluate_transcript_integrity

        transcript.integrity = evaluate_transcript_integrity(
            [item.to_dict() for item in segments], duration_seconds=duration
        )
        identity = hashlib.sha256(url.encode("utf-8")).hexdigest()[:10]
        base = f"{safe_stem(title)}-{identity}"
        artifacts = write_transcript_artifacts(
            transcript, sidecar.parent / "transcript", base, source_name=title
        )
        content_segments = [
            ContentSegment(
                f"speech-{index}", item.start, "speech", text=item.text,
                location={"timestamp_start": item.start, "timestamp_end": item.end},
                metadata={
                    "language": language,
                    "engine": "platform-subtitle",
                    "subtitle_kind": "automatic" if automatic else "manual",
                },
            )
            for index, item in enumerate(segments, 1)
            if item.text
        ]
        return ExtractionResult(
            title=title,
            media_type="video",
            segments=content_segments,
            # The public subtitle is the reproducible source evidence for this
            # branch.  Mark it as the owned source so archival never invents a
            # page.html path that does not exist.
            source_path=sidecar,
            original_uri=url,
            checksum=sha256_file(sidecar),
            transcript_path=artifacts["txt"],
            metadata={
                "platform": platform,
                "platform_id": str(info.get("id") or ""),
                "platform_uploader": str(info.get("uploader") or ""),
                "content_scope": "full_media_transcript",
                "public_access_only": True,
                "cookies_used": False,
                "proxy_used": False,
                "source_subtitle": str(sidecar),
                "subtitle_kind": "automatic" if automatic else "manual",
                "transcription": transcript.metadata(),
            },
        )

    def extract(self, url: str, context: ExtractionContext) -> ExtractionResult:
        if not self.supports(url):
            raise ValueError("不是受支持的公开视频平台链接")
        parsed = urllib.parse.urlsplit(url)
        platform = self.host_platforms[(parsed.hostname or "").casefold()]
        identity = hashlib.sha256(url.encode("utf-8")).hexdigest()[:16]
        platform_cache = context.paths.cache / "public-platform"
        platform_cache.mkdir(parents=True, exist_ok=True)
        directory = context.own_temporary_path(
            Path(tempfile.mkdtemp(prefix=f"{identity}-", dir=platform_cache))
        )
        try:
            return self._extract_in_directory(
                url=url,
                platform=platform,
                identity=identity,
                directory=directory,
                context=context,
            )
        except BaseException as error:
            # Keep the original extraction/cancellation failure authoritative.
            # A cleanup failure is recorded on the context and attached as a
            # diagnostic note so it is neither silent nor allowed to mask it.
            try:
                context.cleanup_temporary_path(directory)
            except OSError as cleanup_error:
                error.add_note(f"附加诊断：{cleanup_error}")
            raise

    def _extract_in_directory(
        self,
        *,
        url: str,
        platform: str,
        identity: str,
        directory: Path,
        context: ExtractionContext,
    ) -> ExtractionResult:
        context.message(f"正在读取 {platform} 公开内容信息（不使用 Cookie 或代理）")
        probe_options = self._options(directory, context)
        probe_options["skip_download"] = True
        info = self._run_ytdlp(url, probe_options, download=False)
        self._enforce_not_live(info)
        context.cancellation.check()
        subtitle = self._subtitle_choice(info)
        if subtitle:
            language, automatic = subtitle
            context.message(f"优先下载公开{'自动' if automatic else '人工'}字幕：{language}")
            subtitle_options = self._options(directory, context)
            subtitle_options.update({
                "skip_download": True,
                "writesubtitles": not automatic,
                "writeautomaticsub": automatic,
                "subtitleslangs": [language],
                "subtitlesformat": "vtt/srt/best",
            })
            try:
                subtitle_info = self._run_ytdlp(url, subtitle_options, download=True)
                self._enforce_not_live(subtitle_info or info)
                sidecars = sorted([*directory.glob("*.vtt"), *directory.glob("*.srt")])
                if sidecars:
                    return self._subtitle_result(
                        url=url,
                        platform=platform,
                        info=subtitle_info or info,
                        sidecar=sidecars[0],
                        language=language,
                        automatic=automatic,
                        context=context,
                    )
            except CancelledError:
                raise
            except RuntimeError:
                context.message("公开字幕不可用，正在改用公开媒体流转写")

        self._enforce_declared_limit(info)
        self._enforce_safe_protocols(info)
        context.message("没有可用公开字幕，正在下载公开媒体流")
        media_options = self._options(directory, context)
        media_options.update({"format": "bestaudio/best", "skip_download": False})
        media_info = self._run_ytdlp(url, media_options, download=True)
        self._enforce_not_live(media_info)
        self._enforce_safe_protocols(media_info)
        self._enforce_declared_limit(media_info)
        context.cancellation.check()
        media_files = sorted(
            path for path in directory.iterdir()
            if (
                path.is_file()
                and not path.is_symlink()
                and path.suffix.casefold() in AudioVideoExtractor.suffixes
            )
        )
        if not media_files:
            raise RuntimeError("平台未返回可处理的公开音视频文件")
        if media_files[0].stat().st_size > self.max_download_bytes:
            raise self._limit_error()
        extracted = AudioVideoExtractor().extract(media_files[0], context)
        extracted.source_path = media_files[0]
        extracted.title = str(media_info.get("title") or info.get("title") or extracted.title).strip()
        extracted.original_uri = url
        extracted.metadata.update({
            "platform": platform,
            "platform_id": str(media_info.get("id") or info.get("id") or ""),
            "platform_uploader": str(media_info.get("uploader") or info.get("uploader") or ""),
            "content_scope": "full_media",
            "public_access_only": True,
            "cookies_used": False,
            "proxy_used": False,
            "temporary_source_owned_by": "ExtractionContext",
            "temporary_source_identity": identity,
            "source_media_bytes": media_files[0].stat().st_size,
        })
        return extracted


class WeixinChannelsExtractor:
    """Resolve public Weixin Channels share metadata and playable streams when exposed."""

    api_url = "https://channels.weixin.qq.com/finder-preview/api/feed/get_feed_info"
    allowed_hosts = {"weixin.qq.com", "www.weixin.qq.com", "channels.weixin.qq.com"}

    @classmethod
    def supports(cls, url: str) -> bool:
        parsed = urllib.parse.urlsplit(url)
        host = (parsed.hostname or "").casefold()
        if parsed.scheme.casefold() not in {"http", "https"} or host not in cls.allowed_hosts:
            return False
        return parsed.path.startswith("/sph/") or parsed.path.startswith("/finder-preview/pages/sph")

    @staticmethod
    def _short_uri(url: str) -> str:
        parsed = urllib.parse.urlsplit(url)
        if parsed.path.startswith("/sph/"):
            candidate = parsed.path.split("/sph/", 1)[1].split("/", 1)[0]
        else:
            candidate = urllib.parse.parse_qs(parsed.query).get("id", [""])[0]
        candidate = urllib.parse.unquote(candidate).strip()
        if not re.fullmatch(r"[A-Za-z0-9_-]{5,128}", candidate):
            raise ValueError("无法识别微信视频号短链编号")
        return candidate

    @staticmethod
    def _playable_url(feed: dict[str, object]) -> str | None:
        for key in ("h264VideoInfo", "h265VideoInfo"):
            value = feed.get(key)
            if isinstance(value, dict) and isinstance(value.get("videoUrl"), str) and value["videoUrl"]:
                return str(value["videoUrl"])
        value = feed.get("videoUrl")
        return str(value) if isinstance(value, str) and value else None

    @staticmethod
    def _title(description: str, author: str) -> str:
        first = re.split(r"[。！？\n]", description, maxsplit=1)[0].strip()
        base = first[:72] or "微信视频号内容"
        return f"{base}｜{author}" if author else base

    def extract(self, url: str, context: ExtractionContext) -> ExtractionResult:
        short_uri = self._short_uri(url)
        preview_url = f"https://channels.weixin.qq.com/finder-preview/pages/sph?id={urllib.parse.quote(short_uri)}"
        payload = json.dumps(
            {"baseReq": {"generalToken": ""}, "shortUri": short_uri},
            ensure_ascii=False,
        ).encode("utf-8")
        request = urllib.request.Request(
            self.api_url,
            data=payload,
            headers={
                "Content-Type": "application/json",
                "Origin": "https://channels.weixin.qq.com",
                "Referer": preview_url,
                "User-Agent": "Mozilla/5.0 AI-Jingjing/1.0",
            },
        )
        context.message("正在解析微信视频号分享链接")
        try:
            with urllib.request.urlopen(request, timeout=45) as response:
                raw = response.read(2 * 1024 * 1024 + 1)
        except urllib.error.HTTPError as exc:
            raise RuntimeError(f"微信视频号链接解析失败（HTTP {exc.code}）") from exc
        except urllib.error.URLError as exc:
            raise RuntimeError("暂时无法连接微信视频号，请稍后重试") from exc
        if len(raw) > 2 * 1024 * 1024:
            raise ValueError("微信视频号返回数据超过安全限制")
        try:
            result = json.loads(raw.decode("utf-8"))
        except (UnicodeError, json.JSONDecodeError) as exc:
            raise RuntimeError("微信视频号返回了无法识别的数据") from exc
        data = result.get("data") if isinstance(result, dict) else None
        data = data if isinstance(data, dict) else {}
        feed = data.get("feedInfo")
        feed = feed if isinstance(feed, dict) else {}
        author_info = data.get("authorInfo")
        author_info = author_info if isinstance(author_info, dict) else {}
        description = str(feed.get("description") or "").strip()
        author = str(author_info.get("nickname") or "").strip()
        error_info = data.get("errMsg")
        if not feed and isinstance(error_info, dict):
            reason = str(error_info.get("title") or "此内容暂时无法读取")
            raise RuntimeError(
                f"{reason}。内容可能需要微信登录、仅好友可见、已删除或链接已过期；"
                "请在微信中保存视频文件后导入。"
            )
        title = self._title(description, author)
        platform_metadata: dict[str, object] = {
            "platform": "weixin_channels",
            "platform_author": author,
            "platform_description": description,
            "platform_cover_url": str(feed.get("coverUrl") or ""),
            "platform_created_at": feed.get("createtime"),
            "weixin_short_uri": short_uri,
        }
        playable_url = self._playable_url(feed)
        if playable_url:
            context.message("已取得视频流，准备下载并转写")
            return DirectMediaURLExtractor().extract_remote(
                playable_url,
                context,
                original_uri=url,
                title=title,
                metadata={**platform_metadata, "content_scope": "full_media"},
                request_headers={"Referer": preview_url},
            )
        raise RuntimeError(
            "微信视频号公开分享页只提供说明和封面，没有开放真实音视频流；"
            "为避免把简介误当成视频内容，本次不入库。请在微信中保存原视频文件后拖入导入。"
        )


def url_extractor_for(url: str):
    if WeixinChannelsExtractor.supports(url):
        return WeixinChannelsExtractor()
    if WeixinArticleExtractor.supports(url):
        return WeixinArticleExtractor()
    if PublicPlatformVideoExtractor.supports(url):
        return PublicPlatformVideoExtractor()
    if DirectMediaURLExtractor.supports(url):
        return DirectMediaURLExtractor()
    return WebExtractor()


FILE_EXTRACTORS = [TextExtractor(), PDFExtractor(), DOCXExtractor(), PPTXExtractor(), ImageExtractor(), AudioVideoExtractor()]


def extractor_for(path: Path):
    suffix = path.suffix.casefold()
    return next((extractor for extractor in FILE_EXTRACTORS if suffix in extractor.suffixes), None)
